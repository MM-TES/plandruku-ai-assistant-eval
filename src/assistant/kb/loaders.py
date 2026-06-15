"""Per-format text extraction for the knowledge base.

Each loader returns a list of ``Segment`` (locator + text) so we keep page/slide/
sheet provenance. Heavy/optional paths (.doc via Word, OCR via Tesseract) are
gated on tool availability and degrade gracefully (raise -> caller skips + logs).
"""
from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from src.utils.logger import setup_logger

_logger = setup_logger(__name__)

TEXT_EXTS = {".docx", ".pdf", ".xlsx", ".xlsm", ".pptx", ".ppt", ".html", ".htm",
             ".txt", ".md", ".csv"}
DOC_EXTS = {".doc"}
IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".tif", ".tiff", ".bmp", ".gif"}

_MAX_CELLS_PER_SHEET = 4000  # guard against huge spreadsheets


@dataclass
class Segment:
    locator: str  # e.g. "стор. 3", "слайд 5", "аркуш Лист1", ""
    text: str


# --- plain text -------------------------------------------------------------
def load_text(path: Path) -> list[Segment]:
    for enc in ("utf-8", "utf-16", "cp1251", "latin-1"):
        try:
            return [Segment("", path.read_text(encoding=enc))]
        except (UnicodeDecodeError, UnicodeError):
            continue
    return [Segment("", path.read_bytes().decode("utf-8", errors="ignore"))]


# --- docx -------------------------------------------------------------------
def load_docx(path: Path) -> list[Segment]:
    from docx import Document

    doc = Document(str(path))
    parts = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text and c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return [Segment("", "\n".join(parts))]


# --- pdf (pymupdf; OCR fallback per empty page) -----------------------------
def load_pdf(path: Path, *, ocr: bool = False, ocr_langs: str = "ukr+rus") -> list[Segment]:
    import fitz  # pymupdf

    from src.assistant import config

    max_ocr = int(config.kb_param("ocr_max_pages", 100))  # bound per-file OCR time
    segments: list[Segment] = []
    ocr_done = 0
    with fitz.open(str(path)) as doc:
        for i, page in enumerate(doc, start=1):
            text = page.get_text("text") or ""
            if not text.strip() and ocr and ocr_done < max_ocr:
                text = _ocr_pdf_page(page)
                ocr_done += 1
            if text.strip():
                segments.append(Segment(f"стор. {i}", text))
    return segments


def _ocr_pdf_page(page) -> str:
    try:
        import io

        import numpy as np
        from PIL import Image

        from src.assistant import config
        from src.assistant.kb import ocr as ocr_engine

        dpi = int(config.kb_param("ocr_dpi", 300))  # higher DPI -> better OCR on scans
        pix = page.get_pixmap(dpi=dpi)
        img = Image.open(io.BytesIO(pix.tobytes("png"))).convert("RGB")
        return ocr_engine.ocr_array(np.array(img))
    except Exception as exc:  # noqa: BLE001
        _logger.debug("PDF OCR skipped: %s", exc)
        return ""


# --- xlsx / xlsm ------------------------------------------------------------
def load_xlsx(path: Path) -> list[Segment]:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    segments: list[Segment] = []
    for ws in wb.worksheets:
        lines: list[str] = []
        n = 0
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None and str(c).strip()]
            if cells:
                lines.append(" | ".join(cells))
                n += len(cells)
            if n > _MAX_CELLS_PER_SHEET:
                break
        if lines:
            segments.append(Segment(f"аркуш {ws.title}", "\n".join(lines)))
    wb.close()
    return segments


# --- pptx -------------------------------------------------------------------
def load_pptx(path: Path) -> list[Segment]:
    from pptx import Presentation

    prs = Presentation(str(path))
    segments: list[Segment] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            note = slide.notes_slide.notes_text_frame.text
            if note.strip():
                parts.append(f"[нотатки] {note}")
        if parts:
            segments.append(Segment(f"слайд {i}", "\n".join(parts)))
    return segments


# --- html -------------------------------------------------------------------
def load_html(path: Path) -> list[Segment]:
    from bs4 import BeautifulSoup

    raw = load_text(path)[0].text
    soup = BeautifulSoup(raw, "html.parser")
    for tag in soup(["script", "style"]):
        tag.decompose()
    text = "\n".join(line.strip() for line in soup.get_text("\n").splitlines() if line.strip())
    return [Segment("", text)]


# --- .doc (legacy) via a single reused MS Word COM instance -----------------
_word_app = None


def _get_word():
    """Open Word once and reuse it (per-file open/quit is slow + leaks processes)."""
    global _word_app
    if _word_app is None:
        import atexit

        import pythoncom
        import win32com.client

        pythoncom.CoInitialize()
        _word_app = win32com.client.DispatchEx("Word.Application")
        _word_app.Visible = False
        _word_app.DisplayAlerts = False
        atexit.register(close_word)
    return _word_app


def close_word() -> None:
    global _word_app
    if _word_app is not None:
        try:
            _word_app.Quit()
        except Exception:  # noqa: BLE001
            pass
        _word_app = None


def load_doc(path: Path) -> list[Segment]:
    word = _get_word()
    doc = word.Documents.Open(str(path), ReadOnly=True, AddToRecentFiles=False)
    try:
        text = doc.Content.Text
    finally:
        doc.Close(False)
    return [Segment("", text)]


# --- image OCR (EasyOCR) ----------------------------------------------------
def load_image(path: Path) -> list[Segment]:
    from src.assistant.kb import ocr as ocr_engine

    text = ocr_engine.ocr_path(path)
    return [Segment("", text)] if text.strip() else []


# --- dispatch ---------------------------------------------------------------
def supported(ext: str, *, include_doc: bool, ocr: bool) -> bool:
    ext = ext.lower()
    if ext in TEXT_EXTS:
        return True
    if include_doc and ext in DOC_EXTS:
        return True
    if ocr and ext in IMAGE_EXTS:
        return True
    return False


def extract(path: Path, *, include_doc: bool = False, ocr: bool = False,
            ocr_langs: str = "ukr+rus") -> list[Segment]:
    """Extract segments from *path*. Raises on failure (caller logs + skips)."""
    ext = path.suffix.lower()
    if ext == ".docx":
        return load_docx(path)
    if ext == ".pdf":
        return load_pdf(path, ocr=ocr, ocr_langs=ocr_langs)
    if ext in (".xlsx", ".xlsm"):
        return load_xlsx(path)
    if ext == ".pptx":
        return load_pptx(path)
    if ext in (".html", ".htm"):
        return load_html(path)
    if ext in (".txt", ".md", ".csv"):
        return load_text(path)
    if ext in DOC_EXTS and include_doc:
        return load_doc(path)
    if ext in IMAGE_EXTS and ocr:
        return load_image(path)
    return []


def ocr_available() -> bool:
    from src.assistant.kb import ocr as ocr_engine

    return ocr_engine.ocr_available()


__all__ = ["Segment", "extract", "supported", "tesseract_available",
           "TEXT_EXTS", "DOC_EXTS", "IMAGE_EXTS"]
